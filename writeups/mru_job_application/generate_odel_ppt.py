#!/usr/bin/env python3
"""
Generate MRU ODeL Strategic Proposal PowerPoint Presentation
Muteesa I Royal University — Branded, Professional, Executive-Ready
Square-cornered design. All costs in USD (1 USD = 3,600 UGX).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os, shutil

# ============================================================
# PALETTE
# ============================================================
NAVY       = RGBColor(0x0A, 0x1F, 0x44)
BLUE       = RGBColor(0x00, 0x3D, 0x7A)
MID_BLUE   = RGBColor(0x1A, 0x5C, 0x9E)
SKY        = RGBColor(0x3A, 0x8F, 0xCF)
GOLD       = RGBColor(0xC8, 0x96, 0x16)
PALE_GOLD  = RGBColor(0xE8, 0xC8, 0x5A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF7, 0xF7, 0xF7)
LGRAY      = RGBColor(0xEE, 0xEE, 0xEE)
MGRAY      = RGBColor(0x88, 0x88, 0x88)
DGRAY      = RGBColor(0x2D, 0x2D, 0x2D)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x1E, 0x8C, 0x52)
ORANGE     = RGBColor(0xD4, 0x7B, 0x1E)

# ============================================================
# LAYOUT
# ============================================================
SW = Inches(13.333); SH = Inches(7.5)
ML = Inches(0.9); CW = Inches(11.533); STRIP = Inches(0.12)
TY = Inches(0.55); TH = Inches(0.75); ABY = Inches(1.25)
BY = Inches(1.65); FY = Inches(7.22); SNY = Inches(7.02)
FN = "Calibri"; FL = "Calibri Light"

prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
TOTAL = 18

# ============================================================
# TRANSITIONS
# ============================================================
def transition(slide, t="fade"):
    sld = slide._element
    for old in sld.findall(qn('p:transition')): sld.remove(old)
    tr = etree.SubElement(sld, qn('p:transition'))
    tr.set('spd', 'med'); tr.set('advClick', '1')
    if t == "fade": etree.SubElement(tr, qn('p:fade'))
    elif t == "push":
        p = etree.SubElement(tr, qn('p:push')); p.set('dir', 'l')
    elif t == "wipe":
        w = etree.SubElement(tr, qn('p:wipe')); w.set('dir', 'r')

# ============================================================
# SHAPE HELPERS (ALL SQUARE CORNERS)
# ============================================================
def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def box(s, l, t, w, h, c):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def txt(s, l, t, w, h, text, sz=18, c=DGRAY, b=False, a=PP_ALIGN.LEFT,
        f=FN, it=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b
    r.font.name = f; r.font.italic = it; tf.paragraphs[0].alignment = a
    return tb

def mtxt(s, l, t, w, h, lines, sz=16, c=DGRAY, b=False, f=FN, a=PP_ALIGN.LEFT, sp=10):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln; r.font.size = Pt(sz)
        r.font.color.rgb = c; r.font.bold = b; r.font.name = f
        p.alignment = a; p.space_after = Pt(sp)
    return tb

def dot(s, l, t, sz, c):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Pt(sz), Pt(sz))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

def circ(s, l, t, d, c):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

# ============================================================
# FURNITURE
# ============================================================
def strip(s, c=NAVY): box(s, Inches(0), Inches(0), STRIP, SH, c)
def footer(s):
    box(s, Inches(0), FY, SW, Pt(3), GOLD)
    txt(s, ML, Inches(7.28), Inches(6), Inches(0.2),
        "Muteesa I Royal University  |  ODeL Strategic Proposal", sz=8, c=MGRAY)
def pgnum(s, n):
    txt(s, Inches(12.0), SNY, Inches(1.0), Inches(0.3),
        f"{n}/{TOTAL}", sz=9, c=MGRAY, a=PP_ALIGN.RIGHT)
def stitle(s, t, c=NAVY):
    txt(s, ML, TY, CW, TH, t, sz=32, c=c, b=True, f=FL)
    box(s, ML, ABY, Inches(1.6), Pt(4), GOLD)

def cslide(title, num, sc=NAVY, bgc=WHITE, tc=NAVY, tr="fade"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, bgc); strip(s, sc); footer(s); pgnum(s, num)
    stitle(s, title, tc); transition(s, tr)
    return s

def secslide(title, sub, num, tr="fade"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    box(s, Inches(0.9), Inches(2.6), Pt(5), Inches(2.2), GOLD)
    txt(s, Inches(1.4), Inches(2.7), Inches(10), Inches(1.2),
        title, sz=40, c=WHITE, b=True, f=FL)
    txt(s, Inches(1.4), Inches(4.1), Inches(10), Inches(0.8),
        sub, sz=20, c=PALE_GOLD, it=True)
    pgnum(s, num); transition(s, tr)
    return s


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY); transition(s, "fade")
box(s, Inches(0), Inches(0), Inches(0.30), SH, GOLD)
box(s, Inches(1.1), Inches(1.2), Inches(10.5), Pt(1), GOLD)
txt(s, Inches(1.1), Inches(1.6), Inches(10), Inches(0.9),
    "Restoring & Reimagining", sz=46, c=WHITE, b=True, f=FL)
txt(s, Inches(1.1), Inches(2.5), Inches(10), Inches(0.9),
    "ODeL at Muteesa I Royal University", sz=46, c=GOLD, b=True, f=FL)
txt(s, Inches(1.1), Inches(3.9), Inches(8), Inches(0.6),
    "A Strategic Proposal to the Vice Chancellor\u2019s Office",
    sz=22, c=RGBColor(0xAA,0xBB,0xCC), f=FL)
box(s, Inches(1.1), Inches(4.8), Inches(2.5), Pt(3), GOLD)
txt(s, Inches(1.1), Inches(5.2), Inches(6), Inches(0.4),
    "Presented by: Muhindo Mubaraka", sz=18, c=WHITE)
txt(s, Inches(1.1), Inches(5.65), Inches(6), Inches(0.4),
    "February 2026", sz=15, c=MGRAY)
txt(s, Inches(1.1), Inches(6.15), Inches(6), Inches(0.4),
    "Confidential \u2014 For Internal Use Only", sz=11,
    c=RGBColor(0x77,0x88,0x99), it=True)
box(s, Inches(1.1), Inches(6.8), Inches(10.5), Pt(1), GOLD)


# ============================================================
# SLIDE 2 — AGENDA  (updated)
# ============================================================
s = cslide("Agenda", 2, tr="push")
agenda = [
    "1.   Understanding ODeL & Its National Significance",
    "2.   Where Uganda\u2019s Universities Stand Today",
    "3.   The Current Crisis at MRU",
    "4.   What We Lose Every Semester",
    "5.   The Opportunity Ahead",
    "6.   Platform Branding & Architecture",
    "7.   Platform Capabilities & Custom Integrations",
    "8.   Why Cloud VPS \u2014 Not a Local Server",
    "9.   Institutional Adoption Policies",
    "10.  Phased Implementation Roadmap",
    "11.  Training & Capacity Building",
    "12.  Revenue Model & Budget",
    "13.  Risks & Mitigation",
    "14.  Immediate Next Steps",
]
mtxt(s, ML, BY + Inches(0.1), Inches(8), Inches(5.5),
     agenda, sz=19, c=DGRAY, sp=11)
box(s, Inches(10.2), BY + Inches(0.3), Inches(2.5), Inches(4.8), NAVY)
txt(s, Inches(10.4), BY + Inches(1.0), Inches(2.1), Inches(1),
    "18", sz=64, c=GOLD, b=True, a=PP_ALIGN.CENTER)
txt(s, Inches(10.4), BY + Inches(2.4), Inches(2.1), Inches(0.5),
    "slides", sz=18, c=WHITE, a=PP_ALIGN.CENTER)
txt(s, Inches(10.4), BY + Inches(3.0), Inches(2.1), Inches(0.5),
    "~20 min", sz=15, c=PALE_GOLD, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 — WHAT IS ODeL
# ============================================================
s = cslide("What is ODeL?", 3)
txt(s, ML, BY, CW, Inches(1.0),
    "Open, Distance and e-Learning (ODeL) is a formally recognised mode of higher "
    "education delivery in Uganda. It enables students who cannot attend face-to-face "
    "classes \u2014 due to distance, employment, or personal commitments \u2014 to earn "
    "accredited qualifications through digital platforms.", sz=17, c=DGRAY)

cards = [
    ("NCHE Mandated", "The National Council for Higher\nEducation now requires all\nuniversities to maintain\nfunctional ODeL systems."),
    ("Post-COVID Priority", "Government authorised and\naccelerated ODeL adoption after\nCOVID-19 exposed critical gaps\nin traditional delivery."),
    ("Strategic Growth Tool", "ODeL unlocks new student\npopulations \u2014 working\nprofessionals, rural communities\nand cross-border learners."),
]
cc = [NAVY, BLUE, MID_BLUE]
for i, (t, d) in enumerate(cards):
    x = ML + Inches(i * 3.95); y = Inches(3.3)
    box(s, x, y, Inches(3.65), Inches(3.3), cc[i])
    txt(s, x+Inches(0.35), y+Inches(0.3), Inches(3.0), Inches(0.5),
        t, sz=19, c=GOLD, b=True)
    box(s, x+Inches(0.35), y+Inches(0.85), Inches(1.2), Pt(2), GOLD)
    txt(s, x+Inches(0.35), y+Inches(1.1), Inches(2.95), Inches(2.0),
        d, sz=15, c=WHITE)


# ============================================================
# SLIDE 4 — NATIONAL LANDSCAPE
# ============================================================
s = cslide("ODeL Across Uganda \u2014 Where Peers Stand", 4)
txt(s, ML, BY, CW, Inches(0.6),
    "Uganda\u2019s leading universities are investing heavily in ODeL. MRU risks falling behind.",
    sz=17, c=MGRAY, it=True)

unis = [
    ("Makerere University", "IODeL institute & MUELE\nplatform. Full ODeL policy\nframework in place.", GREEN, "MATURE"),
    ("Kyambogo University", "ODeL Centre est. 2007.\nAfrican Virtual University\npartnership. Blended learning.", GREEN, "ESTABLISHED"),
    ("Gulu University",     "ODELL Dept launched 2025.\nPilot: B.Ed (Primary).\nTarget: 11,000+ students.", ORANGE, "LAUNCHING"),
    ("KIU",                 "CEODL est. 2011. Diplomas,\nBachelors & Masters via\ndistance learning.", GREEN, "OPERATIONAL"),
    ("Muteesa I Royal (MRU)","Server DOWN for 2 years.\nNo functional ODeL system.\nNo distance learning.", RED, "OFFLINE"),
]
for i, (nm, desc, bc, bt) in enumerate(unis):
    x = Inches(0.55 + i*2.5); y = Inches(2.65)
    bgc = RGBColor(0xFF,0xEB,0xEB) if i==4 else LGRAY
    box(s, x, y, Inches(2.35), Inches(4.0), bgc)
    box(s, x+Inches(0.15), y+Inches(0.2), Inches(2.05), Inches(0.38), bc)
    txt(s, x+Inches(0.15), y+Inches(0.2), Inches(2.05), Inches(0.38),
        bt, sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.15), y+Inches(0.75), Inches(2.05), Inches(0.65),
        nm, sz=13, c=NAVY, b=True)
    txt(s, x+Inches(0.15), y+Inches(1.5), Inches(2.05), Inches(2.2),
        desc, sz=12, c=DGRAY)


# ============================================================
# SLIDE 5 — SECTION: CRISIS
# ============================================================
secslide("The Current Crisis at MRU",
         "Understanding where we stand \u2014 and what it\u2019s costing us", 5)


# ============================================================
# SLIDE 6 — CRISIS DETAILS
# ============================================================
s = cslide("Where MRU Stands Today", 6, sc=RED, tc=RED)
crisis = [
    ("Complete Server Failure",
     "MRU\u2019s ODeL server has been non-functional for over two years. No online courses, "
     "no digital assessments, no distance learning of any kind."),
    ("NCHE Non-Compliance Risk",
     "The National Council for Higher Education requires every university to maintain a "
     "functional ODeL system. Our current state puts MRU at direct regulatory risk."),
    ("Students Turning Elsewhere",
     "Learners who need flexible, distance-based study are enrolling at competing "
     "institutions that offer online programmes."),
    ("Staff Frustration & Low Morale",
     "Faculty have no digital platform for course delivery. Administrative staff lack "
     "tools for student management and reporting."),
    ("Reputational Damage",
     "While peer universities digitise rapidly, MRU\u2019s inability to offer ODeL undermines "
     "our standing among prospective students and partners."),
]
for i, (h, d) in enumerate(crisis):
    y = BY + Inches(i * 1.0)
    dot(s, ML, y+Inches(0.12), 10, RED)
    txt(s, ML+Inches(0.35), y, Inches(3.5), Inches(0.4), h, sz=16, c=RED, b=True)
    txt(s, ML+Inches(0.35), y+Inches(0.38), CW-Inches(0.5), Inches(0.5), d, sz=14, c=DGRAY)


# ============================================================
# SLIDE 7 — COST OF INACTION
# ============================================================
s = cslide("What We Lose Every Semester Without ODeL", 7)
losses = [
    ("Lost Revenue",    "Distance-learning students\nenrol at competing\ninstitutions instead.", "UGX"),
    ("Regulatory Risk",  "Each semester increases\nthe chance of NCHE\nsanctions or conditions.", "!"),
    ("Talent Drain",     "Faculty seek institutions\nwith modern digital\nteaching infrastructure.", "\u2192"),
    ("Missed Funding",   "Donors & partners\nincreasingly require ODeL\nas a prerequisite.", "\u2715"),
]
for i, (t, d, ic) in enumerate(losses):
    x = ML + Inches(i*2.95); y = Inches(2.2)
    circ(s, x+Inches(0.85), y, Inches(0.8), RED)
    txt(s, x+Inches(0.9), y+Inches(0.12), Inches(0.7), Inches(0.55),
        ic, sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s, x, y+Inches(1.1), Inches(2.7), Inches(0.5),
        t, sz=19, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    txt(s, x, y+Inches(1.65), Inches(2.7), Inches(1.5),
        d, sz=14, c=MGRAY, a=PP_ALIGN.CENTER)

box(s, ML, Inches(5.5), CW, Inches(1.0), RGBColor(0xFF,0xF3,0xE0))
txt(s, ML+Inches(0.4), Inches(5.65), CW-Inches(0.8), Inches(0.7),
    '\u201cEvery semester of inaction is not just a missed opportunity \u2014 '
    'it is an active loss of students, credibility, and institutional relevance.\u201d',
    sz=16, c=ORANGE, b=True, a=PP_ALIGN.CENTER, it=True)


# ============================================================
# SLIDE 8 — SECTION: SOLUTION
# ============================================================
secslide("The Proposed Solution",
         "A modern, reliable, and cost-effective ODeL platform for MRU", 8, "push")


# ============================================================
# SLIDE 9 — THE OPPORTUNITY
# ============================================================
s = cslide("What a Working ODeL System Unlocks", 9)
opps = [
    ("Expanded Access",
     "Reach working professionals, rural students, and cross-border "
     "learners who cannot attend in-person classes."),
    ("New Revenue Streams",
     "Distance programmes generate tuition revenue with lower "
     "marginal costs than traditional classroom delivery."),
    ("NCHE Compliance",
     "Meet regulatory requirements and position MRU as a compliant, "
     "forward-thinking institution during evaluations."),
    ("Competitive Positioning",
     "Join Makerere, Kyambogo, and Gulu among Uganda\u2019s digitally "
     "capable universities. Attract students and partnerships."),
]
for i, (t, d) in enumerate(opps):
    x = ML + Inches(i*2.95); y = BY + Inches(0.3)
    box(s, x, y, Inches(2.7), Pt(5), GOLD)
    box(s, x, y+Pt(5), Inches(2.7), Inches(3.8), LGRAY)
    txt(s, x+Inches(0.2), y+Inches(0.25), Inches(2.3), Inches(0.5),
        t, sz=19, c=NAVY, b=True)
    box(s, x+Inches(0.2), y+Inches(0.8), Inches(0.8), Pt(2), GOLD)
    txt(s, x+Inches(0.2), y+Inches(1.05), Inches(2.3), Inches(2.5),
        d, sz=14, c=DGRAY)


# ============================================================
# SLIDE 10 — PLATFORM BRANDING & ARCHITECTURE
# ============================================================
s = cslide("Platform Branding & Architecture", 10, tr="push")

txt(s, ML, BY, CW, Inches(0.5),
    "Three proposed names for MRU\u2019s new ODeL platform \u2014 for the VC\u2019s consideration:",
    sz=16, c=MGRAY, it=True)

brands = [
    ("MRU Learn", "Simple, clean, and institutional.\nInstantly recognisable as\nMRU\u2019s official learning platform."),
    ("RoyalEdge", "Plays on \u201cRoyal\u201d from Muteesa I\nRoyal University. Conveys\nmodern competitive advantage."),
    ("MuteesaLMS", "Ties directly to the university\u2019s\nidentity. Clear purpose:\nMuteesa\u2019s Learning Management System."),
]
for i, (nm, desc) in enumerate(brands):
    x = ML + Inches(i * 3.95); y = Inches(2.35)
    box(s, x, y, Inches(3.65), Inches(0.5), GOLD)
    txt(s, x+Inches(0.2), y+Inches(0.05), Inches(3.2), Inches(0.4),
        nm, sz=20, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    box(s, x, y+Inches(0.5), Inches(3.65), Inches(1.3), LGRAY)
    txt(s, x+Inches(0.2), y+Inches(0.65), Inches(3.2), Inches(1.0),
        desc, sz=14, c=DGRAY, a=PP_ALIGN.CENTER)

# Architecture tiers below
txt(s, ML, Inches(4.15), CW, Inches(0.4),
    "System Architecture", sz=20, c=NAVY, b=True)
box(s, ML, Inches(4.45), Inches(1.6), Pt(3), GOLD)

tiers = [
    ("USERS", "Students  \u00b7  Lecturers  \u00b7  Administrators",
     "Browser or mobile, anywhere with internet.", SKY),
    ("APPLICATION", "Moodle LMS  +  Custom MRU Modules",
     "Courses, assessments, grading, analytics, notifications.", BLUE),
    ("INFRASTRUCTURE", "Cloud VPS: 12 CPU \u00b7 24 GB RAM \u00b7 2 TB SSD",
     "Automated backups, SSL, 99.9% uptime, scalable.", NAVY),
]
for i, (tier, comp, det, clr) in enumerate(tiers):
    y = Inches(4.8 + i * 0.82)
    box(s, Inches(1.2), y, Inches(11.0), Inches(0.7), clr)
    txt(s, Inches(1.5), y+Inches(0.05), Inches(1.5), Inches(0.3),
        tier, sz=10, c=PALE_GOLD, b=True)
    txt(s, Inches(1.5), y+Inches(0.3), Inches(4.5), Inches(0.35),
        comp, sz=14, c=WHITE, b=True)
    txt(s, Inches(7.0), y+Inches(0.15), Inches(4.8), Inches(0.4),
        det, sz=12, c=RGBColor(0xDD,0xEE,0xFF))


# ============================================================
# SLIDE 11 — PLATFORM CAPABILITIES
# ============================================================
s = cslide("What Students & Staff Will Experience", 11)

lf = [
    "Course content upload & management",
    "Assignment submission & online grading",
    "Discussion forums & direct messaging",
    "Video lecture hosting & playback",
    "Timed online quizzes & examinations",
    "Attendance & progress tracking",
]
rf = [
    "Mobile-friendly access (phone & tablet)",
    "Plagiarism detection integration",
    "SMS / WhatsApp notifications",
    "Admin dashboards & enrolment analytics",
    "Automated NCHE-aligned reporting",
    "Student performance trend analysis",
]

box(s, ML, BY+Inches(0.1), Inches(5.5), Inches(0.45), NAVY)
txt(s, ML+Inches(0.2), BY+Inches(0.12), Inches(5), Inches(0.4),
    "Core Moodle Features", sz=16, c=WHITE, b=True)
box(s, Inches(7.0), BY+Inches(0.1), Inches(5.5), Inches(0.45), GOLD)
txt(s, Inches(7.2), BY+Inches(0.12), Inches(5), Inches(0.4),
    "Custom-Built MRU Integrations", sz=16, c=NAVY, b=True)

for i, feat in enumerate(lf):
    y = BY + Inches(0.75 + i*0.62)
    dot(s, ML+Inches(0.1), y+Inches(0.1), 8, BLUE)
    txt(s, ML+Inches(0.45), y, Inches(5), Inches(0.5), feat, sz=15, c=DGRAY)

for i, feat in enumerate(rf):
    y = BY + Inches(0.75 + i*0.62)
    dot(s, Inches(7.1), y+Inches(0.1), 8, GOLD)
    txt(s, Inches(7.45), y, Inches(5), Inches(0.5), feat, sz=15, c=DGRAY)

box(s, Inches(6.65), BY+Inches(0.7), Pt(1.5), Inches(3.8), LGRAY)
box(s, ML, Inches(6.3), CW, Inches(0.5), LGRAY)
txt(s, ML+Inches(0.3), Inches(6.33), CW-Inches(0.6), Inches(0.45),
    "Custom integrations will be developed specifically for MRU\u2019s workflows and NCHE requirements.",
    sz=12, c=MGRAY, a=PP_ALIGN.CENTER, it=True)


# ============================================================
# SLIDE 12 — VPS vs LOCAL SERVER
# ============================================================
s = cslide("Why Cloud VPS \u2014 Not a Local Server", 12)

hy = BY + Inches(0.15); hh = Inches(0.45)
cx = [ML, Inches(4.1), Inches(8.6)]
cw = [Inches(3.1), Inches(4.4), Inches(3.9)]

box(s, cx[0], hy, cw[0], hh, NAVY)
txt(s, cx[0], hy+Inches(0.02), cw[0], hh, "Factor", sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)
box(s, cx[1], hy, cw[1], hh, RED)
txt(s, cx[1], hy+Inches(0.02), cw[1], hh, "Local Server  (Current)", sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)
box(s, cx[2], hy, cw[2], hh, GREEN)
txt(s, cx[2], hy+Inches(0.02), cw[2], hh, "Cloud VPS  (Proposed)", sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)

rows = [
    ("Uptime",            "Power cuts, hardware failure.\nCurrently 0% uptime.",       "99.9% guaranteed.\n24/7 professional monitoring."),
    ("Maintenance",       "On-site IT staff, spare\nparts, cooling systems.",           "Managed by provider.\nMinimal in-house burden."),
    ("Security",          "Vulnerable to physical\ndamage, theft, attacks.",            "Enterprise firewalls, SSL,\nautomated security patches."),
    ("Scalability",       "Fixed capacity. Upgrade\nrequires new hardware.",            "Scale up/down instantly\nbased on student demand."),
    ("Disaster Recovery", "No automated backups.\nHIGH data loss risk.",               "Daily automated backups.\nGeographic redundancy."),
    ("3-Year Cost",       "High: hardware + power\n+ cooling + staff + repairs.",      "Lower TCO. Predictable\npay-as-you-grow pricing."),
]
for i, (fac, loc, cld) in enumerate(rows):
    y = hy + hh + Inches(0.08) + Inches(i*0.72)
    bgc = OFF_WHITE if i%2==0 else WHITE
    for j in range(3): box(s, cx[j], y, cw[j], Inches(0.68), bgc)
    txt(s, cx[0]+Inches(0.15), y+Inches(0.08), cw[0]-Inches(0.3), Inches(0.55),
        fac, sz=13, c=NAVY, b=True)
    txt(s, cx[1]+Inches(0.15), y+Inches(0.05), cw[1]-Inches(0.3), Inches(0.6),
        loc, sz=11, c=DGRAY)
    txt(s, cx[2]+Inches(0.15), y+Inches(0.05), cw[2]-Inches(0.3), Inches(0.6),
        cld, sz=11, c=DGRAY)


# ============================================================
# SLIDE 13 — INSTITUTIONAL ADOPTION POLICIES  (NEW)
# ============================================================
s = cslide("Ensuring Adoption \u2014 Institutional Policies", 13)

txt(s, ML, BY, CW, Inches(0.6),
    "Technology succeeds only when people use it. The following realistic, phased policies "
    "will drive adoption across MRU without being unreasonably rigid.",
    sz=17, c=MGRAY, it=True)

policies = [
    ("Compulsory Assignment Submission",
     "All coursework assignments must be submitted through the ODeL platform. "
     "This ensures every student and lecturer interacts with the system regularly "
     "and creates a verifiable digital record of academic work."),
    ("Lecture Materials Upload Requirement",
     "Lecturers must upload course outlines, reading materials, and lecture notes "
     "to the platform at the start of each semester. This makes the system the "
     "single source of truth for all course content."),
    ("Online Grade Publication",
     "All continuous assessment marks and final grades must be published through "
     "the platform. Students check results online, reducing administrative "
     "bottlenecks and paper-based processes."),
    ("Student Course Registration via Platform",
     "Course registration and enrolment for each semester should be processed "
     "through the ODeL system, ensuring all students have active accounts "
     "and are familiar with the platform from day one."),
    ("Blended Learning Attendance Tracking",
     "For programmes using blended delivery, attendance for online sessions "
     "must be tracked via the platform\u2019s built-in tools, ensuring accountability "
     "and enabling data-driven monitoring of student engagement."),
]

for i, (title, desc) in enumerate(policies):
    y = BY + Inches(0.7) + Inches(i * 0.97)
    # Number badge
    box(s, ML, y, Inches(0.4), Inches(0.4), NAVY)
    txt(s, ML, y+Inches(0.01), Inches(0.4), Inches(0.38),
        str(i+1), sz=16, c=GOLD, b=True, a=PP_ALIGN.CENTER)
    # Title
    txt(s, ML+Inches(0.55), y, Inches(3.5), Inches(0.4),
        title, sz=15, c=NAVY, b=True)
    # Description
    txt(s, ML+Inches(0.55), y+Inches(0.35), CW-Inches(0.7), Inches(0.55),
        desc, sz=12, c=DGRAY)


# ============================================================
# SLIDE 14 — ROADMAP
# ============================================================
s = cslide("Implementation Roadmap \u2014 3 Phases", 14, tr="push")

phases = [
    ("PHASE 1", "Foundation & Pilot", "Months 1\u20133",
     ["Procure & configure Cloud VPS",
      "Install & customise Moodle LMS",
      "Develop custom MRU integrations",
      "Pilot with 3\u20135 selected courses",
      "Gather feedback & iterate"], SKY),
    ("PHASE 2", "Training & Expansion", "Months 4\u20136",
     ["Faculty & staff training workshops",
      "Student onboarding & orientation",
      "Expand to 15\u201320 courses",
      "Establish IT helpdesk & support",
      "Refine based on pilot learnings"], BLUE),
    ("PHASE 3", "Full Rollout", "Months 7\u201312",
     ["University-wide deployment",
      "All programmes available on ODeL",
      "Performance monitoring & reporting",
      "NCHE compliance documentation",
      "Continuous improvement cycle"], NAVY),
]
for i, (ph, tl, tm, items, clr) in enumerate(phases):
    x = ML + Inches(i*3.95); y = BY + Inches(0.2)
    box(s, x, y, Inches(3.65), Inches(0.95), clr)
    txt(s, x+Inches(0.25), y+Inches(0.05), Inches(1.5), Inches(0.3),
        ph, sz=12, c=PALE_GOLD, b=True)
    txt(s, x+Inches(0.25), y+Inches(0.35), Inches(2.5), Inches(0.45),
        tl, sz=18, c=WHITE, b=True)
    txt(s, x+Inches(2.3), y+Inches(0.05), Inches(1.1), Inches(0.3),
        tm, sz=11, c=PALE_GOLD, a=PP_ALIGN.RIGHT)
    for j, item in enumerate(items):
        iy = y + Inches(1.15 + j*0.62)
        dot(s, x+Inches(0.15), iy+Inches(0.1), 7, clr)
        txt(s, x+Inches(0.5), iy, Inches(3.0), Inches(0.5), item, sz=14, c=DGRAY)
    if i < 2:
        txt(s, x+Inches(3.7), y+Inches(0.2), Inches(0.3), Inches(0.5),
            "\u2192", sz=26, c=GOLD, b=True, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 — TRAINING
# ============================================================
s = cslide("Training & Capacity Building", 15)
txt(s, ML, BY, CW, Inches(0.6),
    "Technology alone is not enough. Successful ODeL requires people who are confident "
    "and capable using it.", sz=17, c=MGRAY, it=True)

groups = [
    ("Faculty & Lecturers", [
        "Hands-on Moodle workshops",
        "Online pedagogy best practices",
        "Ongoing helpdesk & documentation"], NAVY),
    ("Students", [
        "Orientation sessions at enrolment",
        "Step-by-step video tutorials",
        "SMS/WhatsApp support channel"], BLUE),
    ("IT & Admin Staff", [
        "System administration training",
        "Monitoring & backup procedures",
        "NCHE reporting workflows"], MID_BLUE),
]
for i, (grp, items, clr) in enumerate(groups):
    x = ML + Inches(i*3.95); y = Inches(2.65)
    box(s, x, y, Inches(3.65), Inches(0.5), clr)
    txt(s, x+Inches(0.2), y+Inches(0.05), Inches(3.2), Inches(0.4),
        grp, sz=17, c=WHITE, b=True)
    for j, item in enumerate(items):
        iy = y + Inches(0.75 + j*0.75)
        dot(s, x+Inches(0.15), iy+Inches(0.12), 7, clr)
        txt(s, x+Inches(0.5), iy, Inches(3.0), Inches(0.6), item, sz=14, c=DGRAY)

box(s, ML, Inches(5.7), CW, Inches(0.9), RGBColor(0xE8,0xF0,0xF8))
txt(s, ML+Inches(0.4), Inches(5.85), CW-Inches(0.8), Inches(0.6),
    "Training will be scheduled during existing faculty development periods to minimise "
    "disruption. All materials will also be available online for self-paced review.",
    sz=13, c=BLUE, a=PP_ALIGN.CENTER, it=True)


# ============================================================
# SLIDE 16 — REVENUE MODEL & BUDGET  (NEW / UPDATED)
# ============================================================
s = cslide("Revenue Model & Budget", 16)

# --- REVENUE SECTION ---
txt(s, ML, BY, Inches(5.5), Inches(0.4),
    "Revenue: Technology Fee (Precedent: Kyambogo University)", sz=18, c=NAVY, b=True)
box(s, ML, BY+Inches(0.4), Inches(5.5), Pt(3), GOLD)

rev_lines = [
    "Technology fee per student:   UGX 5,000 / semester",
    "Estimated active students:    1,000",
    "",
    "Revenue per semester:         UGX 5,000,000",
    "Revenue per year (2 sem):     UGX 10,000,000",
    "",
    "Converted at 1 USD = 3,600 UGX:",
    "Annual revenue:               $2,778",
]
mtxt(s, ML, BY+Inches(0.55), Inches(5.5), Inches(3.2),
     rev_lines, sz=14, c=DGRAY, sp=5)

# Revenue highlight box
box(s, ML, Inches(5.1), Inches(5.5), Inches(0.55), RGBColor(0xE8,0xF5,0xE9))
txt(s, ML+Inches(0.2), Inches(5.15), Inches(5.1), Inches(0.45),
    "\u2713  Annual tech fee revenue ($2,778) exceeds total platform cost ($2,500)",
    sz=14, c=GREEN, b=True)

# --- BUDGET SECTION ---
txt(s, Inches(7.0), BY, Inches(5.5), Inches(0.4),
    "Budget Breakdown (Annual, in USD)", sz=18, c=NAVY, b=True)
box(s, Inches(7.0), BY+Inches(0.4), Inches(5.5), Pt(3), GOLD)

# Budget table
bx = Inches(7.0); bw1 = Inches(3.3); bw2 = Inches(2.2)
# Header
box(s, bx, BY+Inches(0.55), bw1, Inches(0.38), NAVY)
txt(s, bx+Inches(0.1), BY+Inches(0.57), bw1-Inches(0.2), Inches(0.34),
    "Item", sz=12, c=WHITE, b=True)
box(s, bx+bw1, BY+Inches(0.55), bw2, Inches(0.38), NAVY)
txt(s, bx+bw1+Inches(0.1), BY+Inches(0.57), bw2-Inches(0.2), Inches(0.34),
    "Cost (USD)", sz=12, c=WHITE, b=True, a=PP_ALIGN.RIGHT)

budget_items = [
    ("Cloud VPS Hosting", "12 CPU \u00b7 24 GB RAM \u00b7 2 TB SSD", "$1,000 / year"),
    ("System Development,\nDeployment & Support", "Moodle setup, custom modules,\nongoing technical support", "$1,000 / year"),
    ("Training & Capacity Building", "Faculty, student & admin\ntraining workshops + materials", "$500"),
]

for i, (item, desc, cost) in enumerate(budget_items):
    y = BY + Inches(0.98) + Inches(i * 0.9)
    bgc = OFF_WHITE if i % 2 == 0 else WHITE
    box(s, bx, y, bw1, Inches(0.85), bgc)
    box(s, bx+bw1, y, bw2, Inches(0.85), bgc)
    txt(s, bx+Inches(0.1), y+Inches(0.05), bw1-Inches(0.2), Inches(0.35),
        item, sz=12, c=NAVY, b=True)
    txt(s, bx+Inches(0.1), y+Inches(0.4), bw1-Inches(0.2), Inches(0.4),
        desc, sz=10, c=MGRAY)
    txt(s, bx+bw1+Inches(0.1), y+Inches(0.15), bw2-Inches(0.2), Inches(0.5),
        cost, sz=14, c=DGRAY, b=True, a=PP_ALIGN.RIGHT)

# Total row
ty = BY + Inches(0.98) + Inches(3 * 0.9)
box(s, bx, ty, bw1+bw2, Inches(0.45), NAVY)
txt(s, bx+Inches(0.1), ty+Inches(0.05), bw1-Inches(0.2), Inches(0.35),
    "TOTAL ANNUAL COST", sz=13, c=GOLD, b=True)
txt(s, bx+bw1+Inches(0.1), ty+Inches(0.05), bw2-Inches(0.2), Inches(0.35),
    "$2,500", sz=16, c=WHITE, b=True, a=PP_ALIGN.RIGHT)

# Bottom note
box(s, ML, Inches(6.3), CW, Inches(0.55), RGBColor(0xE8,0xF5,0xE9))
txt(s, ML+Inches(0.2), Inches(6.35), CW-Inches(0.4), Inches(0.45),
    "The project is fully self-sustaining: a modest UGX 5,000 technology fee per student "
    "per semester covers all costs with a $278 annual surplus.",
    sz=12, c=GREEN, b=True, a=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 17 — RISKS
# ============================================================
s = cslide("Risks & How We\u2019ll Manage Them", 17)
risks = [
    ("Limited Internet",
     "Some students have poor connectivity.",
     "Mobile-optimised platform, offline content downloads, SMS notifications, low-bandwidth mode."),
    ("Adoption Resistance",
     "Faculty or students may resist digital tools.",
     "Hands-on training, peer champions programme, phased rollout with early adopters."),
    ("Budget Constraints",
     "Limited funds for upfront investment.",
     "Self-funding via tech fee, open-source tools, clear ROI projections, phased spending."),
    ("Long-Term Sustainability",
     "Concern about maintenance after setup.",
     "Knowledge transfer to IT staff, full documentation, automated monitoring, maintenance plan."),
]
for i, (risk, issue, mit) in enumerate(risks):
    y = BY + Inches(i * 1.3)
    box(s, ML, y, Inches(2.8), Inches(0.4), ORANGE)
    txt(s, ML+Inches(0.15), y+Inches(0.02), Inches(2.5), Inches(0.35),
        risk, sz=13, c=WHITE, b=True)
    txt(s, ML, y+Inches(0.45), Inches(5.3), Inches(0.5),
        issue, sz=13, c=MGRAY, it=True)
    txt(s, Inches(6.2), y+Inches(0.15), Inches(0.5), Inches(0.4),
        "\u2192", sz=22, c=GREEN, b=True, a=PP_ALIGN.CENTER)
    box(s, Inches(6.7), y, Inches(5.7), Inches(1.05), RGBColor(0xE8,0xF5,0xE9))
    txt(s, Inches(6.9), y+Inches(0.15), Inches(5.3), Inches(0.8),
        mit, sz=12, c=DGRAY)


# ============================================================
# SLIDE 18 — NEXT STEPS & CLOSE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY); transition(s, "fade")
box(s, Inches(0), Inches(0), Inches(0.30), SH, GOLD)
box(s, Inches(1.1), Inches(0.5), Inches(10.5), Pt(1), GOLD)
txt(s, Inches(1.1), Inches(0.8), Inches(8), TH,
    "Immediate Next Steps", sz=36, c=WHITE, b=True, f=FL)
box(s, Inches(1.1), Inches(1.5), Inches(1.8), Pt(3), GOLD)

steps = [
    ("1", "Approval to proceed with the ODeL restoration project"),
    ("2", "Form a small ODeL Task Force \u2014 IT, Faculty, Admin representation"),
    ("3", "Approve the UGX 5,000 technology fee to fund the project"),
    ("4", "Procure Cloud VPS hosting and professional domain"),
    ("5", "Begin Phase 1 pilot \u2014 first courses online within 90 days"),
    ("6", "Schedule monthly progress reviews with VC\u2019s office"),
]
for i, (num, step) in enumerate(steps):
    y = Inches(1.85 + i * 0.72)
    circ(s, Inches(1.1), y, Inches(0.45), GOLD)
    txt(s, Inches(1.15), y+Inches(0.04), Inches(0.38), Inches(0.38),
        num, sz=18, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    txt(s, Inches(1.75), y+Inches(0.05), Inches(9.5), Inches(0.4),
        step, sz=17, c=WHITE)

box(s, Inches(1.1), Inches(6.1), Inches(10.5), Pt(1), GOLD)
txt(s, Inches(1.1), Inches(6.3), Inches(10.5), Inches(0.7),
    '\u201cThis is not just about fixing a server. It is about positioning '
    'Muteesa I Royal University for the future of higher education in Uganda.\u201d',
    sz=16, c=PALE_GOLD, a=PP_ALIGN.CENTER, it=True)
pgnum(s, 18)


# ============================================================
# SAVE
# ============================================================
out_dir = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(out_dir, "MRU_ODeL_Strategic_Proposal.pptx")
desk = os.path.expanduser("~/Desktop/MRU_ODeL_Strategic_Proposal.pptx")
prs.save(out); shutil.copy2(out, desk)
print(f"\u2705 Saved: {out}")
print(f"\u2705 Desktop: {desk}")
print(f"   Slides: {len(prs.slides)}")
